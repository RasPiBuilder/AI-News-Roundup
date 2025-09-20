# ai_news_roundup.py

# === Core Python ===
import os
import random
import requests
import re
from datetime import datetime
import glob
import gc
import io
import logging
import traceback
from pathlib import Path
import tempfile

# === Third-party ===
from ddgs import DDGS
from groq import Groq
import pyttsx3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, UnidentifiedImageError
import win32com.client  # windows only
from moviepy.editor import (
    ImageClip,
    AudioFileClip,
    VideoFileClip,
    concatenate_videoclips,
    vfx,
    afx,
)

# ---------------- CONFIG ---------------- #

TOPICS = {
    "Anthropic": ["Anthropic news", "Claude AI updates", "Anthropic announcements"],
    "OpenAI": ["OpenAI news", "OpenAI announcements", "GPT-5 updates", "ChatGPT news"],
    "Humanoid Robots": ["humanoid robot news", "robotics breakthroughs", "AI-powered robots"],
    "AI Startups": ["AI startup venture capital", "AI funding news", "AI startup acquisitions"],
}

SITES = [
    "theverge.com",
    "arstechnica.com",
    "techcrunch.com",
    "the-decoder.com",
]

SEARCH_MODIFIERS = [
    "September 2025",
    "latest",
    "breaking",
    "update",
    "research",
]

OUTPUT_DIR = "output"
AUDIO_DIR = os.path.join(OUTPUT_DIR, "audio_clips")
IMAGE_DIR = os.path.join(OUTPUT_DIR, "images")
SLIDE_IMG_DIR = os.path.join(OUTPUT_DIR, "slides")
CLIPS_DIR = os.path.join(OUTPUT_DIR, "clips")
PPTX_FILE = os.path.join(OUTPUT_DIR, "news_roundup.pptx")
VIDEO_FILE = os.path.join(OUTPUT_DIR, "news_roundup.mp4")
FINAL_AUDIO = os.path.join(OUTPUT_DIR, "news_roundup.mp3")
LOG_FILE = os.path.join(OUTPUT_DIR, "build.log")

# Ensure dirs exist
for d in (OUTPUT_DIR, AUDIO_DIR, IMAGE_DIR, SLIDE_IMG_DIR, CLIPS_DIR):
    os.makedirs(d, exist_ok=True)

# ---------------- LOGGING ---------------- #

def _setup_logging():
    logger = logging.getLogger("roundup")
    logger.setLevel(logging.DEBUG)
    # Clear handlers if re-run
    logger.handlers = []
    fmt = logging.Formatter(
        "%(asctime)s | %(levelname)-8s | %(name)s | %(message)s",
        datefmt="%Y-%m-%d %H:%M:%S",
    )

    # File handler (DEBUG+)
    fh = logging.FileHandler(LOG_FILE, encoding="utf-8")
    fh.setLevel(logging.DEBUG)
    fh.setFormatter(fmt)
    logger.addHandler(fh)

    # Console handler (INFO+)
    ch = logging.StreamHandler()
    ch.setLevel(logging.INFO)
    ch.setFormatter(fmt)
    logger.addHandler(ch)

    logger.debug("Logging initialized.")
    return logger

log = _setup_logging()

# ---------------- GROQ SETUP ---------------- #
# (Keep API key for rapid testing as requested)
GROQ_API_KEY = "YOUR-API-KEY-HERE"
if not GROQ_API_KEY:
    log.warning("GROQ_API_KEY not set. Set it via environment variable before running.")
client = Groq(api_key=GROQ_API_KEY) if GROQ_API_KEY else None

def groq_call(prompt, max_tokens=300):
    if client is None:
        raise RuntimeError("GROQ client not configured (missing GROQ_API_KEY).")
    try:
        response = client.chat.completions.create(
            model="meta-llama/llama-4-scout-17b-16e-instruct",
            messages=[
                {"role": "system", "content": "Output ONLY what is asked for."},
                {"role": "user", "content": prompt},
            ],
            temperature=0.2,
            max_completion_tokens=max_tokens,
        )
        out = response.choices[0].message.content.strip()
        log.debug("GROQ call succeeded (len=%d).", len(out))
        return out
    except Exception as e:
        log.error("GROQ call failed: %s\n%s", e, traceback.format_exc())
        raise

# ---------------- HELPERS: IMAGES ---------------- #

def _sanitize_keywords(keywords_text):
    parts = [p.strip("•-–— \t\"'`") for p in re.split(r"[,\n;/]+", keywords_text) if p.strip()]
    return ", ".join(parts[:2]) if parts else ""

def validate_image_bytes(raw: bytes, min_size=100) -> bool:
    try:
        with Image.open(io.BytesIO(raw)) as im:
            im.verify()
        with Image.open(io.BytesIO(raw)) as im:
            im.load()
            w, h = im.size
        return w >= min_size and h >= min_size
    except Exception:
        return False

def _has_alpha_channel(im: Image.Image) -> bool:
    return im.mode in ("LA", "RGBA", "PA") or (im.mode == "P" and "transparency" in im.info)

def convert_image_bytes_to_png_path(raw: bytes, out_path_no_ext: str) -> str:
    """Convert arbitrary image bytes to PNG on disk and return the path."""
    png_path = out_path_no_ext + ".png"
    with Image.open(io.BytesIO(raw)) as im:
        im = im.convert("RGBA")
        if not _has_alpha_channel(im):
            im = im.convert("RGB")
        im.save(png_path, format="PNG", optimize=True)
    return png_path

def fetch_image(keywords, filename_stub):
    """
    Fetch an image via DDGS, ensure it's an image, validate, convert to PNG, and return PNG path.
    Returns None if nothing valid found.
    """
    kw = _sanitize_keywords(keywords)
    if not kw:
        log.info("No keywords for image fetch.")
        return None

    log.info("Searching images for: %s", kw)
    with DDGS() as ddgs:
        for r in ddgs.images(kw, max_results=8):
            url = r.get("image") or r.get("thumbnail")
            if not url:
                continue
            try:
                resp = requests.get(url, timeout=12)
                if resp.status_code != 200:
                    continue

                ctype = resp.headers.get("Content-Type", "")
                if ctype and not ctype.lower().startswith("image/"):
                    continue

                raw = resp.content
                if not validate_image_bytes(raw):
                    continue

                out_no_ext = os.path.join(IMAGE_DIR, filename_stub)
                png_path = convert_image_bytes_to_png_path(raw, out_no_ext)

                # Disk validation
                try:
                    with Image.open(png_path) as im:
                        im.verify()
                except Exception:
                    try:
                        os.remove(png_path)
                    except Exception:
                        pass
                    continue

                log.info("Image saved: %s", png_path)
                return png_path

            except Exception as e:
                log.debug("Image fetch attempt failed: %s", e)
                continue

    log.warning("No valid image found for keywords: %s", kw)
    return None

# ---------------- HELPERS: AUDIO ---------------- #

def save_audio(text, filename):
    try:
        engine = pyttsx3.init()
        voices = engine.getProperty("voices")
        for v in voices:
            if "Zira" in v.name:
                engine.setProperty("voice", v.id)
                break
        engine.setProperty("rate", 175)
        engine.setProperty("volume", 0.9)
        engine.save_to_file(text, filename)
        engine.runAndWait()
        log.info("Queued audio: %s (len=%d chars)", filename, len(text))
    except Exception as e:
        log.error("Audio synthesis failed for %s: %s", filename, e)
        raise

# ---------------- AI TEXT ---------------- #

def search_snippets(query, num_results=5):
    snippets = []
    with DDGS() as ddgs:
        for r in ddgs.text(query, max_results=num_results):
            title = r.get("title", "").strip()
            body = r.get("body", "").strip()
            if title or body:
                snippets.append(f"{title}: {body}".strip(": "))
    log.debug("Found %d snippets for query: %s", len(snippets), query)
    return snippets

def get_bullet_points(summary):
    prompt = (
        "Turn the following news summary into 5–7 concise bullet points.\n"
        "- Output ONLY the bullet lines, one per line.\n"
        "- No numbering, no symbols, no headings.\n"
        "- Keep each line under 18 words.\n\n"
        f"{summary}"
    )
    return groq_call(prompt)

def get_script(summary, bullets):
    prompt = (
        "Write a natural narration (~220–350 words) using the material below.\n"
        "- Output ONLY the narration text (plain paragraphs).\n\n"
        f"SUMMARY:\n{summary}\n\n"
        f"BULLETS:\n{bullets}"
    )
    return groq_call(prompt, max_tokens=600)

def get_image_keywords(summary):
    prompt = (
        "From the summary, output 1–2 short keyword phrases for an image search.\n"
        "- Output ONLY the keywords, comma-separated if two.\n\n"
        f"{summary}"
    )
    return groq_call(prompt, max_tokens=30)

def get_intro_text(date_str, topics):
    topic_list = ", ".join(topics)
    prompt = (
        f"Write a short intro (2–3 sentences) for a news roundup on {date_str}. "
        f"Mention: {topic_list}. Output ONLY the text."
    )
    return groq_call(prompt, max_tokens=160)

def get_outro_text():
    prompt = "Write a short outro (1–2 sentences). Output ONLY the text."
    return groq_call(prompt, max_tokens=120)

# ---------------- PPT BUILD (non-overlapping layout) ---------------- #

def _add_bullets_to_frame(tf, bullets_str):
    """Fill a text_frame with bullet lines, with mild autosize."""
    lines = [ln.strip() for ln in bullets_str.splitlines() if ln.strip()]
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    if not lines:
        p0 = tf.paragraphs[0]
        p0.text = ""
        return

    # Font sizing heuristic
    n = len(lines)
    if n <= 6:
        size = 18
    elif n <= 10:
        size = 16
    else:
        size = 14

    # First bullet
    p0 = tf.paragraphs[0]
    p0.text = lines[0]
    p0.level = 0
    p0.font.size = Pt(size)

    for ln in lines[1:]:
        p = tf.add_paragraph()
        p.text = ln
        p.level = 0
        p.font.size = Pt(size)

def _add_image_fit_center(slide, image_path, left_in, top_in, max_w_in, max_h_in):
    """Add an image scaled to fit within the given box (no overlap)."""
    with Image.open(image_path) as im:
        w, h = im.size
    max_w_px = int(max_w_in * 96)
    max_h_px = int(max_h_in * 96)
    scale = min(max_w_px / w, max_h_px / h, 1.0)
    disp_w_in = (w * scale) / 96.0
    disp_h_in = (h * scale) / 96.0
    left = Inches(left_in + (max_w_in - disp_w_in) / 2)
    top = Inches(top_in + (max_h_in - disp_h_in) / 2)
    return slide.shapes.add_picture(image_path, left, top, width=Inches(disp_w_in))

def build_ppt(segments):
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]  # Title + Content
    title_slide_layout = prs.slide_layouts[0]

    # Intro slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "AI & Tech News Roundup"
    slide.placeholders[1].text = datetime.now().strftime("%B %d, %Y")
    subtitle = slide.placeholders[1].text_frame.paragraphs[0]
    subtitle.font.size = Pt(24)
    subtitle.alignment = PP_ALIGN.CENTER

    # Topic slides (non-overlapping two-column layout)
    for seg in segments:
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = seg["topic"]

        # LEFT column: bullets
        body = slide.placeholders[1]
        body.left = Inches(0.6)
        body.top = Inches(1.4)
        body.width = Inches(6.0)
        body.height = Inches(4.6)

        tf = body.text_frame
        _add_bullets_to_frame(tf, seg["bullets"])

        # RIGHT column: image
        if seg.get("image") and os.path.exists(seg["image"]):
            try:
                _add_image_fit_center(
                    slide,
                    seg["image"],
                    left_in=6.8,
                    top_in=1.4,
                    max_w_in=3.2,
                    max_h_in=4.6,
                )
            except UnidentifiedImageError:
                log.warning("Could not identify image for topic '%s'; skipping.", seg["topic"])
            except Exception as e:
                log.error("Error placing image for '%s': %s", seg["topic"], e)
        else:
            log.info("No valid image for topic '%s'", seg["topic"])

    # Outro slide
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Thanks for Watching"
    slide.placeholders[1].text = "Stay tuned for tomorrow's update!"
    outro_para = slide.placeholders[1].text_frame.paragraphs[0]
    outro_para.font.size = Pt(24)
    outro_para.alignment = PP_ALIGN.CENTER

    prs.save(PPTX_FILE)
    log.info("Saved presentation: %s", PPTX_FILE)

# ---------------- SLIDE EXPORT W/ STANDARDIZED NAMES ---------------- #

def export_slides_to_images(pptx_file, out_dir=SLIDE_IMG_DIR):
    """
    Deterministic export: we explicitly export each slide to the exact filename we want.
    - intro_slide.jpg
    - topic_01_slide.jpg ... topic_NN_slide.jpg
    - outro_slide.jpg
    This avoids any reliance on PowerPoint's default 'Slide1.JPG' naming.
    """
    pptx_file = os.path.abspath(pptx_file)
    out_dir = os.path.abspath(out_dir)

    # Clear old standardized files
    for pat in ("intro_slide.*", "topic_*.jpg", "topic_*.png", "outro_slide.*", "Slide*.JPG", "Slide*.PNG"):
        for f in glob.glob(os.path.join(out_dir, pat)):
            try:
                os.remove(f)
            except Exception:
                pass

    log.info("Exporting slides to images via PowerPoint (per-slide export)…")
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1

    # Important: WithWindow=False keeps it headless/fast.
    pres = powerpoint.Presentations.Open(pptx_file, WithWindow=False)
    try:
        total = pres.Slides.Count
        if total < 2:
            raise RuntimeError(f"Presentation has {total} slide(s); need at least intro + outro.")

        # Map slides deterministically:
        # 1 -> intro, 2..(total-1) -> topics, total -> outro
        for idx in range(1, total + 1):
            slide = pres.Slides(idx)
            if idx == 1:
                fname = "intro_slide.jpg"
            elif idx == total:
                fname = "outro_slide.jpg"
            else:
                topic_idx = idx - 1  # topic 1 is slide 2
                fname = f"topic_{topic_idx:02}_slide.jpg"

            dst = os.path.join(out_dir, fname)
            slide.Export(dst, "JPG")
            if not os.path.exists(dst):
                raise RuntimeError(f"PowerPoint failed to export slide {idx} to {dst}")

            log.debug("Exported: slide %d -> %s", idx, dst)

    finally:
        pres.Close()
        powerpoint.Quit()

    # Build return dict
    topics = []
    # Topic slides are from slide 2..(total-1) => topic_01..topic_{total-2}
    for i in range(1, total - 1):
        topics.append(os.path.join(out_dir, f"topic_{i:02}_slide.jpg"))

    result = {
        "intro": os.path.join(out_dir, "intro_slide.jpg"),
        "topics": topics,
        "outro": os.path.join(out_dir, "outro_slide.jpg"),
        "ext": ".jpg",
    }
    log.info("Slides standardized: %s, %d topics, %s", result["intro"], len(topics), result["outro"])
    return result


# ---------------- VIDEO BUILDING (clip-per-section, then stitch) ---------------- #

def _image_with_audio_clip(image_path, audio_path, pre_roll=0.6, fade=0.4, tail_pad=0.2):
    """
    Build a clip:
      - Image shows immediately
      - Audio starts after pre_roll
      - Gentle fade-in/out on video and audio
      - Tail pad ensures fade-out isn't cut early
    """
    audio = AudioFileClip(audio_path)
    total_dur = pre_roll + audio.duration + tail_pad

    # Base image clip with fades
    img = ImageClip(image_path).set_duration(total_dur)
    if fade > 0:
        img = img.fx(vfx.fadein, fade).fx(vfx.fadeout, fade)

    # Audio with start offset and fades
    a = audio.set_start(pre_roll)
    if fade > 0:
        a = a.fx(afx.audio_fadein, fade).fx(afx.audio_fadeout, min(fade, max(0.1, tail_pad)))

    return img.set_audio(a)

def _write_clip(clip, out_path, fps=24):
    """Write a single clip to disk, safely closing resources."""
    clip.write_videofile(
        out_path,
        fps=fps,
        codec="libx264",
        audio_codec="aac",
        temp_audiofile=str(Path(out_path).with_suffix(".m4a")),
        remove_temp=True,
        verbose=False,
        logger=None,
    )
    clip.close()

def build_video(segments, intro_audio, outro_audio, pre_roll_seconds=0.6, fade_seconds=0.4, fps=24):
    """
    1) Export slide images and standardize filenames (.jpg or .png preserved).
    2) Build INDIVIDUAL CLIP FILES:
       - 00_intro.mp4
       - 01_topic_XX.mp4 (for each segment)
       - 99_outro.mp4
    3) Concatenate (with per-clip fade in/out already applied).
    4) Extract combined audio to FINAL_AUDIO.
    """
    slides = export_slides_to_images(PPTX_FILE)
    intro_slide = slides["intro"]
    topic_slides = slides["topics"]
    outro_slide = slides["outro"]

    if len(topic_slides) != len(segments):
        log.warning("Topic slide count (%d) != segment count (%d). Will pair by min length.",
                    len(topic_slides), len(segments))

    n = min(len(topic_slides), len(segments))
    if n == 0:
        raise RuntimeError("No topic slides available to build video.")

    clip_paths = []

    # Intro
    intro_path = os.path.join(CLIPS_DIR, "00_intro.mp4")
    log.info("Rendering intro clip -> %s", intro_path)
    intro_clip = _image_with_audio_clip(
        intro_slide, intro_audio, pre_roll=pre_roll_seconds, fade=fade_seconds
    )
    _write_clip(intro_clip, intro_path, fps=fps)
    clip_paths.append(intro_path)

    # Topics 1..n
    for idx in range(1, n + 1):
        seg = segments[idx - 1]
        slide_path = topic_slides[idx - 1]
        topic_clip_out = os.path.join(CLIPS_DIR, f"01_topic_{idx:02}.mp4")
        log.info("Rendering topic %02d '%s' -> %s", idx, seg["topic"], topic_clip_out)
        try:
            clip = _image_with_audio_clip(
                slide_path, seg["audio"], pre_roll=pre_roll_seconds, fade=fade_seconds
            )
            _write_clip(clip, topic_clip_out, fps=fps)
            clip_paths.append(topic_clip_out)
        except Exception as e:
            log.error("Failed rendering topic %02d: %s", idx, e)
            raise

    # Outro
    outro_path = os.path.join(CLIPS_DIR, "99_outro.mp4")
    log.info("Rendering outro clip -> %s", outro_path)
    outro_clip = _image_with_audio_clip(
        outro_slide, outro_audio, pre_roll=pre_roll_seconds, fade=fade_seconds
    )
    _write_clip(outro_clip, outro_path, fps=fps)
    clip_paths.append(outro_path)

    # Stitch
    log.info("Stitching %d clips into final video…", len(clip_paths))
    stitched = concatenate_videoclips([VideoFileClip(p) for p in clip_paths], method="compose")
    stitched.write_videofile(
        VIDEO_FILE,
        fps=fps,
        codec="libx264",
        audio_codec="aac",
        temp_audiofile=str(Path(VIDEO_FILE).with_suffix(".m4a")),
        remove_temp=True,
        verbose=False,
        logger=None,
    )

    # Export combined audio
    with AudioFileClip(VIDEO_FILE) as a:
        a.write_audiofile(FINAL_AUDIO, fps=44100, verbose=False, logger=None)

    stitched.close()
    log.info("Saved video: %s", VIDEO_FILE)
    log.info("Saved audio: %s", FINAL_AUDIO)

# ---------------- MAIN EXECUTION ---------------- #

def main():
    today = datetime.now().strftime("%A, %B %d, %Y")
    log.info("Fetching news for %s", today)

    all_segments = []
    topic_index = 1

    for topic, base_terms in TOPICS.items():
        chosen_base = random.choice(base_terms)
        chosen_sites = random.sample(SITES, k=min(2, len(SITES)))
        site_summaries = []

        log.info("Topic: %s", topic)
        for site in chosen_sites:
            modifier = random.choice(SEARCH_MODIFIERS)
            query = f"{chosen_base} {modifier} site:{site}"
            log.info("  Searching: %s", query)
            snippets = search_snippets(query)
            if not snippets:
                continue
            combined = " ".join(snippets[:5])
            site_summaries.append(f"{site}: {combined[:400]}")

        if not site_summaries:
            log.info("Skipping topic '%s' (no summaries).", topic)
            continue

        try:
            raw_summary = " ".join(site_summaries)
            bullets = get_bullet_points(raw_summary)
            script = get_script(raw_summary, bullets)
            keywords = get_image_keywords(raw_summary)
            img_file = fetch_image(keywords, f"{topic.replace(' ', '_')}")
        except Exception as e:
            log.error("AI processing failed for '%s': %s", topic, e)
            continue

        # Uniform audio naming: topic_XX_audio.wav
        audio_file = os.path.join(AUDIO_DIR, f"topic_{topic_index:02}_audio.wav")
        save_audio(script, audio_file)
        topic_index += 1

        all_segments.append(
            {
                "topic": topic,
                "bullets": bullets,
                "script": script,
                "audio": audio_file,  # standardized
                "image": img_file,    # PNG path or None (for placing on slide)
            }
        )

    topics_covered = [seg["topic"] for seg in all_segments]
    if topics_covered:
        intro_text = get_intro_text(datetime.now().strftime("%B %d, %Y"), topics_covered)
    else:
        intro_text = f"Welcome to your AI and tech news roundup for {datetime.now().strftime('%B %d, %Y')}."
    outro_text = get_outro_text()

    # Uniform intro/outro audio names
    intro_audio = os.path.join(AUDIO_DIR, "intro_audio.wav")
    outro_audio = os.path.join(AUDIO_DIR, "outro_audio.wav")
    save_audio(intro_text, intro_audio)
    save_audio(outro_text, outro_audio)
    log.info("Saved %d topic clips + intro/outro to %s", len(all_segments), AUDIO_DIR)

    try:
        # Build the deck (bullets + images)
        build_ppt(all_segments)
        # Build final video from per-section clips (precise alignment + fade transitions)
        build_video(all_segments, intro_audio, outro_audio, pre_roll_seconds=0.6, fade_seconds=0.4, fps=24)
    finally:
        gc.collect()
        log.info("All tasks complete, resources cleaned up.")

if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        log.critical("Fatal error: %s\n%s", e, traceback.format_exc())
        raise


